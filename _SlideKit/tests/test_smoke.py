# -*- coding: utf-8 -*-
"""_SlideKit 冒烟测试：折行、spec 校验、母版结构、页码字段、auto-fit、报告页型。
运行： python -m pytest tests/            （需 pytest）
或：   python tests/test_smoke.py         （无 pytest 也能跑）
"""
import os
import sys
import zipfile
import tempfile

import pytest

_SK = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, _SK)
import deck as D          # noqa: E402
import theme as T         # noqa: E402

GOOD = [
    {"kind": "cover", "title": "T", "sub": "s", "src": "J.C"},
    {"kind": "bullets", "title": "要点页", "accent": "1F3A8A", "bullets": ["完整的一句话要点。", "第二条要点。"]},
    {"kind": "table", "title": "对比", "table": {
        "headers": ["项", "值"], "col_align": ["left", "right"], "rows": [["A", "1"], ["B", "2"]]}},
    {"kind": "chart", "title": "趋势", "chart": {
        "type": "line", "categories": ["0", "1", "2"], "series": [("s", [1, 2, 3])]}},
    {"kind": "refs", "title": "参考文献", "refs": ["Author, Title, Venue, 2020."]},
    {"kind": "close", "title": "谢谢", "src": "J.C"},
]


def _build(tmp):
    out = os.path.join(tmp, "smoke.pptx")
    D.build_pptx(GOOD, out, len(GOOD), asset_dir=tmp, page_label="smoke")
    return out


# ------------------------- 折行（不拆词） ------------------------------------ #
def test_wrap_no_midword_split():
    for word in ("supercalifragilistic", "placement", "Innovus"):
        assert D._wrap(word, 5) == [word]            # 整词不拆
        assert T._wrap_w(word, 0.2, 14) == [word]


def test_wrap_keeps_all_tokens():
    out = D._wrap("mesh router placement engine", 8)
    joined = "".join(out)
    for tok in ("mesh", "router", "placement", "engine"):
        assert tok in joined


def test_circ_fallback():
    assert D._circ(0) == "①" and D._circ(9) == "⑩"
    assert D._circ(10) == "(11)" and D._circ(25) == "(26)"   # 超 10 退化，避免越界/缺字形


# ------------------------- spec 校验 ---------------------------------------- #
def test_validate_ok():
    assert D.validate_specs(GOOD) is True


def test_validate_missing_field():
    with pytest.raises(ValueError):
        D.validate_specs([{"kind": "split", "title": "x"}])       # 缺 figure / bullets


def test_validate_unknown_kind():
    with pytest.raises(ValueError):
        D.validate_specs([{"kind": "bogus", "title": "x"}])


def test_validate_table_substructure():
    with pytest.raises(ValueError):
        D.validate_specs([{"kind": "table", "title": "x", "table": {"headers": ["a"]}}])  # 缺 rows


# ------------------------- 母版 / 结构不变量 --------------------------------- #
def test_build_structure():
    from pptx import Presentation
    tmp = tempfile.mkdtemp()
    out = _build(tmp)
    prs = Presentation(out)
    assert len(list(prs.slides._sldIdLst)) == len(GOOD)
    assert len(prs.slide_masters) == 1
    assert len(prs.slide_layouts) == 1
    assert list(prs.slide_layouts[0].placeholders) == []     # 干净版式，无默认占位符


def test_slidenum_field_present():
    tmp = tempfile.mkdtemp()
    out = _build(tmp)
    z = zipfile.ZipFile(out)
    # 第 2 页是内容页（bullets），应带原生 slidenum 字段
    xml = z.read("ppt/slides/slide2.xml").decode("utf-8")
    assert 'type="slidenum"' in xml


def test_body_has_autofit():
    tmp = tempfile.mkdtemp()
    out = _build(tmp)
    xml = zipfile.ZipFile(out).read("ppt/slides/slide2.xml").decode("utf-8")
    assert "normAutofit" in xml                              # 正文框 shrink-to-fit 已开


def test_report_kinds_build_and_reload():
    from pptx import Presentation
    tmp = tempfile.mkdtemp()
    out = _build(tmp)
    prs = Presentation(out)                                  # 重新打开即验证 XML 良构
    assert any(sh.has_table for sh in prs.slides[2].shapes)  # 表格页
    assert any(sh.has_chart for sh in prs.slides[3].shapes)  # 图表页


def test_font_override_and_restore():
    tmp = tempfile.mkdtemp()
    before = (D.YAHEI, D.MONO)
    out = os.path.join(tmp, "f.pptx")
    D.build_pptx(GOOD, out, len(GOOD), asset_dir=tmp, page_label="f", font="Source Han Sans SC")
    assert (D.YAHEI, D.MONO) == before                       # 构建后全局字体还原
    assert "Source Han Sans SC" in zipfile.ZipFile(out).read("ppt/slides/slide2.xml").decode("utf-8")


if __name__ == "__main__":
    fns = [v for k, v in sorted(globals().items()) if k.startswith("test_") and callable(v)]
    for fn in fns:
        fn()
        print("PASS", fn.__name__)
    print("ALL PASS (%d)" % len(fns))
