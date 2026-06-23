# -*- coding: utf-8 -*-
"""
_SlideKit · deck.py  (v3 · 可编辑 PPT)
=====================================
从一份「幻灯片规格」(specs) 同时产出：
  build_pptx(specs)      → 可编辑 .pptx（原生标题/要点文本 + 插入的论文插图）
  build_previews(specs)  → 每页一张 matplotlib 版面预览 PNG（供肉眼校对，因为本机无法渲染 pptx）
两者用同一套版面坐标，预览看着对，pptx 版面就对。

规格 spec = {kind, title, sub, bullets, bullets2, figure, accent, ...}
kind ∈ title | split | cards2 | grid6 | cols2 | close
"""
import os
import sys
import copy
import uuid
import logging

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.image as mpimg
from matplotlib.patches import Rectangle, FancyBboxPatch, Circle, FancyArrowPatch
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import theme as T  # noqa: E402

SLIDE_W, SLIDE_H = 13.333, 7.5
YAHEI = "Microsoft YaHei"
MONO = "Consolas"

# 统一配色（IEEE 靛蓝学术）——单一真源在 theme.py，这里只 lstrip('#') 复用，杜绝三处各写一份 hex。
# 改配色只动 theme.py 即全局生效（含图、deck、消费侧 slides.py）。
PRIMARY   = T.BLUE.lstrip("#")    # 主色 靛蓝：标题条/页码/主卡片/封面竖条
PRIMARY_D = T.BLUE_D.lstrip("#")
ACCENT    = T.AMBER.lstrip("#")   # 暖橙：仅作次列/强调
INK_C     = T.INK.lstrip("#")     # 主文字
SUB_C     = T.INK2.lstrip("#")    # 次文字
MUTED_C   = T.MUTED.lstrip("#")   # 弱化
PAGE_C    = "9AA3B0"              # 页码/页脚（deck 专用浅灰，theme 无对应）
HAIR_C    = T.LINE.lstrip("#")    # 细分隔线
BG_C      = T.WHITE.lstrip("#")   # 母版底色（浅）
PAGE_LABEL = "Floorplan · 布图规划"

# 数据系列配色（图表/对比用）——同样派生自 theme，靛蓝起头
SERIES_COLORS = [c.lstrip("#") for c in (T.BLUE, T.AMBER, T.TEAL, T.ROSE, T.VIOLET)]

_log = logging.getLogger("slidekit")


def _rgb(hexc):
    return RGBColor.from_string(hexc.lstrip("#"))


def fit(box, iw, ih):
    """图片等比放入 box=(x,top,w,h)，返回居中后的 (x,top,w,h)。"""
    x, t, w, h = box
    ar = iw / float(ih)
    if ar > w / float(h):
        nw, nh = w, w / ar
    else:
        nh, nw = h, h * ar
    return x + (w - nw) / 2.0, t + (h - nh) / 2.0, nw, nh


# ============================ pptx 端 ===================================== #
def _set_font(run, name=None, size=14, bold=False, color=INK_C, mono=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    fn = MONO if mono else (name or YAHEI)   # 运行时读取，可被 build_pptx 的 font/mono_font 覆盖
    run.font.name = fn
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", fn)


def _txt(slide, box, lines, sizes=None, colors=None, bolds=None, monos=None,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6, wrap=True, line_sp=None, auto_fit=False):
    x, t, w, h = box
    tb = slide.shapes.add_textbox(Inches(x), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    if auto_fit:                      # 正文框：内容过满时让 PowerPoint 自动缩字号塞下（shrink-to-fit）
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    n = len(lines)
    sizes = sizes or [14] * n
    colors = colors or [INK_C] * n
    bolds = bolds or [False] * n
    monos = monos or [False] * n
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if line_sp:
            p.line_spacing = line_sp
        r = p.add_run(); r.text = ln
        _set_font(r, size=sizes[i], color=colors[i], bold=bolds[i], mono=monos[i])
    return tb


def _bar(slide, x, t, w, h, color):
    sp = slide.shapes.add_shape(1, Inches(x), Inches(t), Inches(w), Inches(h))  # 1=rect
    sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _card(slide, x, t, w, h, fill="FFFFFF", line="CBD5E1", accent=None):
    sp = slide.shapes.add_shape(5, Inches(x), Inches(t), Inches(w), Inches(h))  # 5=rounded rect
    sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    sp.line.color.rgb = _rgb(line); sp.line.width = Pt(1)
    sp.shadow.inherit = False
    if accent:
        _bar(slide, x, t + 0.12, 0.09, h - 0.24, accent)
    return sp


def _title_block(slide, title, sub, accent=PRIMARY, chapter=None):
    # 左侧主色竖条 + 全宽细分隔线已下放到版式（母版）上，这里只填每页不同的标题文本
    _txt(slide, (0.85, 0.42, 9.2, 0.8), [title], sizes=[28], bolds=[True], colors=[INK_C])
    if sub:
        _txt(slide, (0.87, 1.18, 9.2, 0.5), [sub], sizes=[14], colors=[SUB_C])
    if chapter:                              # 右上角章节面包屑（现代结构感）
        _txt(slide, (9.35, 0.55, 3.35, 0.4), [chapter], sizes=[12], bolds=[True], colors=[PRIMARY], align=PP_ALIGN.RIGHT)


def _page(slide):
    """底部右侧：插入 PowerPoint 原生【幻灯片编号字段】(slidenum)，随增删/重排页自动更新；
    比手写 'N/24' 更省事——加页不用改总数。页脚标签在版式上，二者不重叠（左/右分置）。"""
    tb = slide.shapes.add_textbox(Inches(10.8), Inches(7.05), Inches(2.2), Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    fld = p._p.makeelement(qn("a:fld"), {"id": "{%s}" % str(uuid.uuid4()).upper(), "type": "slidenum"})
    rPr = fld.makeelement(qn("a:rPr"), {"lang": "en-US", "sz": "1000"})
    sf = rPr.makeelement(qn("a:solidFill"), {}); cl = sf.makeelement(qn("a:srgbClr"), {"val": PAGE_C})
    sf.append(cl); rPr.append(sf)
    rPr.append(rPr.makeelement(qn("a:latin"), {"typeface": YAHEI}))
    t = fld.makeelement(qn("a:t"), {}); t.text = "1"
    fld.append(rPr); fld.append(t)
    p._p.append(fld)
    return tb


def _layout_chrome(prs, lay, page_label):
    """把每页重复出现的"母版件"放到【版式】上——内容/导览页自动继承，新增页自动带样式；
    要全局改样式（条色/分隔线/页脚）只改这一处，或直接在 PowerPoint 的版式视图里改。
    python-pptx 不支持直接给版式加形状，故先在草稿页上画好，再把形状 XML 复制进版式、删草稿页。
    页脚标签 page_label 随每篇笔记而变（如 'Floorplan · 布图规划' / 'Placement · 布局'），由调用方传入。"""
    scratch = prs.slides.add_slide(lay)
    before = len(scratch.shapes._spTree)
    _bar(scratch, 0.6, 0.5, 0.16, 0.62, PRIMARY)        # 左侧主色竖条
    _bar(scratch, 0.85, 1.6, 11.85, 0.022, HAIR_C)      # 全宽细分隔线
    _txt(scratch, (0.6, 7.06, 6.5, 0.4), [page_label], sizes=[10], colors=[PAGE_C])  # 左下页脚标签（随文档变）
    lay_spTree = lay.shapes._spTree
    for el in list(scratch.shapes._spTree)[before:]:
        lay_spTree.append(copy.deepcopy(el))
    sldIdLst = prs.slides._sldIdLst                      # 删掉这张临时草稿页
    sid = sldIdLst[-1]
    rId = sid.get(qn("r:id"))
    sldIdLst.remove(sid)
    try:
        prs.part.drop_rel(rId)
    except Exception as e:                  # 草稿页关系已不存在等——不致命，但记录便于排查
        _log.debug("drop_rel(%s) 跳过: %s", rId, e)


def _mask(slide):
    """封面/收尾页：盖一层满版白底，遮住版式继承来的内容页母版件，再画自己的整版设计。"""
    _bar(slide, 0, 0, SLIDE_W, SLIDE_H, BG_C)


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(color)


def _missing_pic(slide, path, box):
    """缺图占位：画一只灰色虚线框 + 居中红字 'missing: fXX.png'，并打印告警。
    避免某张图未生成 / 改名 / 新桶引用未画图时，整份 deck 直接 FileNotFoundError 崩掉。"""
    x, t, w, h = box
    sp = slide.shapes.add_shape(5, Inches(x), Inches(t), Inches(w), Inches(h))  # 5=rounded rect
    sp.fill.solid(); sp.fill.fore_color.rgb = _rgb("EDF0F4")
    sp.line.color.rgb = _rgb("C7CDD6"); sp.line.width = Pt(1.2)
    sp.shadow.inherit = False
    _txt(slide, (x, t + h / 2 - 0.35, w, 0.7), ["缺图 missing:", os.path.basename(path)],
         sizes=[14, 13], bolds=[True, False], colors=["A6473C", "6B7280"],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=2)
    print("[warn] 缺图（已用占位框）:", path)


def _pic(slide, path, box):
    if not os.path.isfile(path):
        _missing_pic(slide, path, box)
        return
    iw, ih = Image.open(path).size
    x, t, w, h = fit(box, iw, ih)
    slide.shapes.add_picture(path, Inches(x), Inches(t), Inches(w), Inches(h))


def _bullets_lines(bullets):
    return ["▪  " + b for b in bullets]


CIRC = "①②③④⑤⑥⑦⑧⑨⑩⑪⑫⑬⑭"


def _circ(i):
    """第 i 个圈号（i 从 0 起）。①–⑩ 用圈号字形；超过 10 退化为 (n) 文本——
    既避免 IndexError，也绕开 YaHei 缺 ⑪–⑭ 字形导致的预览/pptx 豆腐块。"""
    return CIRC[i] if i < 10 else f"({i + 1})"


def _marks(items, style="bullet"):
    """并列内容用 ▪；递进/有序内容用 ①②③（style='num'）；叙述性内容不分点（style='prose'，无标记）。"""
    if style == "num":
        return [f"{_circ(i)}  {b}" for i, b in enumerate(items)]
    if style == "prose":
        return list(items)
    return ["▪  " + b for b in items]


def _circ_num(slide, x, y, num, d=0.46, fill=PRIMARY, fg="FFFFFF", fs=15):
    sp = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(d), Inches(d))  # 9=oval
    sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    sp.line.fill.background(); sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num)
    _set_font(r, size=fs, bold=True, color=fg)


def _cover(slide, s):
    """封面：浅底 + 左侧主色厚竖条 + 大标题，独立于内容页版式。"""
    _bg(slide, BG_C)
    _mask(slide)
    _bar(slide, 0.0, 0.0, 0.30, SLIDE_H, PRIMARY)
    _txt(slide, (1.15, 1.45, 11, 0.4), [s.get("tag", "数字 IC 后端 · DVD Lecture 6")],
         sizes=[14], bolds=[True], colors=[PRIMARY])
    _txt(slide, (1.1, 2.45, 11.6, 1.3), [s["title"]], sizes=[46], bolds=[True], colors=[INK_C])
    if s.get("sub"):
        _txt(slide, (1.15, 3.95, 11, 0.6), [s["sub"]], sizes=[19], colors=[SUB_C])
    _bar(slide, 1.15, 4.80, 3.8, 0.045, PRIMARY)
    if s.get("line"):
        _txt(slide, (1.15, 5.08, 11.4, 0.6), [s["line"]], sizes=[14], bolds=[True], colors=[ACCENT])
    _txt(slide, (1.15, 6.85, 11.6, 0.4), [s.get("src", "")], sizes=[11], colors=[MUTED_C])


def _agenda(slide, s):
    """导览：浅底 + 标题条 + 两列编号小节 + 底部主线高亮，独立于内容页版式。"""
    _bg(slide, BG_C)
    _title_block(slide, s["title"], s.get("sub"))
    secs = s["sections"]; colx = [0.95, 7.0]; w = 5.45
    n = len(secs); per = max(1, (n + 1) // 2)   # 两列均分；行数随小节数自适应，避免列索引越界
    top0 = 2.05; rowh = min(1.16, 4.15 / per)   # 小节多时压缩行距，防纵向溢出
    for j, (num, ttl, det) in enumerate(secs):
        x = colx[j // per]; y = top0 + (j % per) * rowh
        _circ_num(slide, x, y, num)
        _txt(slide, (x + 0.64, y - 0.06, w, 0.45), [ttl], sizes=[16], bolds=[True], colors=[INK_C])
        _txt(slide, (x + 0.64, y + 0.40, w, 0.55), [det], sizes=[12.5], colors=[MUTED_C])
    if s.get("line"):
        _card(slide, 0.95, 6.28, 11.0, 0.6, fill="F6E8D5", line="E4C79A")
        _txt(slide, (1.25, 6.41, 10.6, 0.4), [s["line"]], sizes=[13], bolds=[True], colors=["8A5212"])


def _close(slide, s):
    """收尾：浅底 + 左侧主色厚竖条，与封面呼应。"""
    _bg(slide, BG_C)
    _mask(slide)
    _bar(slide, 0.0, 0.0, 0.30, SLIDE_H, PRIMARY)
    _txt(slide, (1.1, 2.75, 11.6, 1.2), [s["title"]], sizes=[44], bolds=[True], colors=[INK_C])
    if s.get("sub"):
        _txt(slide, (1.15, 4.15, 11, 0.6), [s["sub"]], sizes=[18], colors=[SUB_C])
    _bar(slide, 1.15, 4.95, 3.4, 0.045, ACCENT)
    if s.get("line"):
        _txt(slide, (1.15, 5.25, 11.4, 0.5), [s["line"]], sizes=[13], colors=[MUTED_C])
    _txt(slide, (1.15, 6.85, 11.6, 0.4), [s.get("src", "")], sizes=[11], colors=[MUTED_C])


def _section(slide, s):
    """分节标题页（现代分隔）：左侧主色面板 + 大号章节数字，右侧章节标题 + 小节清单。"""
    _bg(slide, BG_C)
    _mask(slide)
    pw = 4.5
    _bar(slide, 0.0, 0.0, pw, SLIDE_H, PRIMARY)
    _txt(slide, (0.0, 1.55, pw, 0.5), [s.get("eyebrow", "CHAPTER")], sizes=[16], bolds=[True],
         colors=["C5D6F7"], align=PP_ALIGN.CENTER)
    _txt(slide, (0.0, 2.2, pw, 2.5), [str(s["num"])], sizes=[150], bolds=[True], colors=["FFFFFF"],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide, (5.1, 1.95, 7.7, 1.2), [s["title"]], sizes=[36], bolds=[True], colors=[INK_C])
    if s.get("sub"):
        _txt(slide, (5.15, 3.2, 7.6, 0.6), [s["sub"]], sizes=[16], colors=[SUB_C])
    _bar(slide, 5.15, 3.95, 2.4, 0.05, ACCENT)
    items = s.get("items", [])
    dy = min(0.52, 2.9 / max(1, len(items)))
    for i, it in enumerate(items):
        y = 4.35 + i * dy
        _bar(slide, 5.15, y + 0.06, 0.13, 0.13, ACCENT)
        _txt(slide, (5.45, y - 0.04, 7.3, 0.4), [it], sizes=[13.5], colors=[INK_C])


def _clean_master(prs, keep):
    """清掉默认模板自带的占位符（日期/页脚/页码/标题/正文）并删掉未使用的版式：
    否则默认占位符会与自定义母版件重叠（日期框压住页脚标签），且默认模板是 4:3、占位符按
    4:3 排布，在 16:9 母版视图里显得不对。只保留实际使用的那张空白版式（已 16:9、放了母版件）。"""
    for m in prs.slide_masters:
        for ph in list(m.placeholders):                       # 母版上的占位符
            ph._element.getparent().remove(ph._element)
        for lay in list(m.slide_layouts):
            if lay._element is keep._element:                 # 保留并清干净这张要用的版式
                for ph in list(lay.placeholders):
                    ph._element.getparent().remove(ph._element)
            else:                                             # 删掉其余未用版式（4:3）
                try:
                    m.slide_layouts.remove(lay)
                except Exception:
                    for ph in list(lay.placeholders):
                        ph._element.getparent().remove(ph._element)


# ---------- 面向技术报告的页型：table / chart / refs（原生可编辑） ---------- #
def _align(a):
    return {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}.get(a, PP_ALIGN.LEFT)


def _cell_border(cell, color=HAIR_C, width_pt=0.75):
    """给表格单元四边描细线（python-pptx 不直接支持，需写 a:tcPr/a:lnL..lnB）。
    四个 ln 元素须排在 fill 之前，故统一 insert(0,…) 插到 tcPr 最前。"""
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):   # 逆序插到最前 → 最终顺序 L,R,T,B
        old = tcPr.find(qn(tag))
        if old is not None:
            tcPr.remove(old)
        ln = tcPr.makeelement(qn(tag), {"w": str(int(width_pt * 12700)), "cap": "flat", "cmpd": "sng", "algn": "ctr"})
        sf = ln.makeelement(qn("a:solidFill"), {})
        clr = sf.makeelement(qn("a:srgbClr"), {"val": color})
        sf.append(clr); ln.append(sf)
        tcPr.insert(0, ln)


def _style_cell(cell, text, size=12, color=INK_C, bold=False, fill="FFFFFF", align="left", font=None):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Pt(7)
    cell.margin_top = cell.margin_bottom = Pt(3)
    cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(fill)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = _align(align)
    r = p.add_run(); r.text = text
    _set_font(r, name=font, size=size, bold=bold, color=color)
    _cell_border(cell)


def _table(slide, s):
    """原生可编辑表格：主色表头(白字粗体) + 斑马纹正文 + 细线边框 + 数值列右对齐。
    spec: {kind:'table', table:{headers:[...], rows:[[...],...], col_align?:[l/r/c], col_w?:[..]}}"""
    t = s["table"]; headers = t["headers"]; rows = t["rows"]
    aligns = t.get("col_align"); widths = t.get("col_w")
    ncol = len(headers); nrow = len(rows) + 1
    x, top, w, h = 0.7, 1.78, 11.95, 5.05
    gtbl = slide.shapes.add_table(nrow, ncol, Inches(x), Inches(top), Inches(w), Inches(h))
    tbl = gtbl.table
    tbl.first_row = False; tbl.horz_banding = False   # 关掉内置样式条纹，改用我们显式的斑马纹
    if widths:
        tot = float(sum(widths))
        for c, cw in enumerate(widths):
            tbl.columns[c].width = Inches(w * cw / tot)
    for r in range(nrow):
        tbl.rows[r].height = Inches(h / nrow)
    al = lambda c: (aligns[c] if aligns and c < len(aligns) else "left")
    for c, htext in enumerate(headers):
        _style_cell(tbl.cell(0, c), str(htext), size=12.5, color="FFFFFF", bold=True, fill=PRIMARY, align=al(c))
    for ri, row in enumerate(rows, 1):
        fill = "FFFFFF" if ri % 2 else "EDF0F4"
        for c in range(ncol):
            _style_cell(tbl.cell(ri, c), str(row[c]) if c < len(row) else "", size=12, fill=fill, align=al(c))


def _chart(slide, s):
    """原生可编辑图表（line / bar）。spec: {kind:'chart', chart:{type:'line'|'bar',
    categories:[...], series:[(name,[v,...]),...], x_title?, y_title?}}"""
    c = s["chart"]; ctype = c.get("type", "line")
    cd = CategoryChartData(); cd.categories = c["categories"]
    for name, vals in c["series"]:
        cd.add_series(name, vals)
    xlc = XL_CHART_TYPE.LINE_MARKERS if ctype == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED
    x, top, w, h = 0.7, 1.78, 11.95, 5.05
    ch = slide.shapes.add_chart(xlc, Inches(x), Inches(top), Inches(w), Inches(h), cd).chart
    ch.has_legend = len(c["series"]) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout = False
    ch.font.size = Pt(11); ch.font.name = YAHEI
    for i, ser in enumerate(ch.series):
        col = _rgb(SERIES_COLORS[i % len(SERIES_COLORS)])
        if ctype == "line":
            ser.format.line.color.rgb = col; ser.format.line.width = Pt(2.0)
        else:
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb = col
    for axis, key in ((ch.category_axis, "x_title"), (ch.value_axis, "y_title")):
        if c.get(key):
            axis.has_title = True; axis.axis_title.text_frame.text = c[key]


def _refs(slide, s):
    """编号参考文献页。spec: {kind:'refs', refs:['Author, Title, Venue, Year.', ...]}"""
    items = s["refs"]
    lines = [f"[{i + 1}]  {r}" for i, r in enumerate(items)]
    if len(items) > 8:
        mid = (len(items) + 1) // 2
        _txt(slide, (0.7, 1.78, 6.0, 5.2), lines[:mid], sizes=[12] * mid, colors=[INK_C] * mid, space=8, line_sp=1.18, auto_fit=True)
        _txt(slide, (6.9, 1.78, 6.0, 5.2), lines[mid:], sizes=[12] * (len(lines) - mid),
             colors=[INK_C] * (len(lines) - mid), space=8, line_sp=1.18, auto_fit=True)
    else:
        _txt(slide, (0.7, 1.78, 12.0, 5.2), lines, sizes=[12.5] * len(lines), colors=[INK_C] * len(lines), space=9, line_sp=1.2, auto_fit=True)


_REQUIRED = {
    "cover": ["title"], "title": ["title"], "close": ["title"],
    "agenda": ["title", "sections"],
    "split": ["title", "figure", "bullets"],
    "bullets": ["title", "bullets"],
    "cols2": ["title", "cols"], "cards2": ["title", "cards"], "grid6": ["title", "items"],
    "table": ["title", "table"], "chart": ["title", "chart"], "refs": ["title", "refs"],
    "section": ["title", "num"],
}


def validate_specs(specs, asset_dir=None):
    """逐页校验：kind 合法、必需字段齐全、table/chart 子结构完整。有硬错误一次列全并抛 ValueError，
    避免到深处才以晦涩的 KeyError 崩溃。asset_dir 给定时顺带检查 split 缺图（只告警，构建用占位框）。"""
    problems, warns = [], []
    for i, s in enumerate(specs, 1):
        k = s.get("kind")
        if k not in _REQUIRED:
            problems.append(f"第{i}页：未知 kind={k!r}（合法：{sorted(_REQUIRED)}）"); continue
        for key in _REQUIRED[k]:
            if key not in s:
                problems.append(f"第{i}页 kind={k}：缺字段 {key!r}")
        if k == "table" and isinstance(s.get("table"), dict):
            if "headers" not in s["table"] or "rows" not in s["table"]:
                problems.append(f"第{i}页 table：table 需含 headers 与 rows")
        if k == "chart" and isinstance(s.get("chart"), dict):
            if "categories" not in s["chart"] or "series" not in s["chart"]:
                problems.append(f"第{i}页 chart：chart 需含 categories 与 series")
        if k == "split" and asset_dir and "figure" in s:
            if not os.path.isfile(os.path.join(asset_dir, s["figure"])):
                warns.append(f"第{i}页：缺图 {s['figure']}（将用占位框）")
    for w in warns:
        print("[warn]", w)
    if problems:
        raise ValueError("spec 校验失败：\n  - " + "\n  - ".join(problems))
    return True


# ---------- 原生框图引擎：一份布局 → PowerPoint 原生形状 + matplotlib 预览（可编辑、不插图片）-------- #
_DROLE = {  # role -> (浅填充, 描边/主色, 文字)  —— 派生自 theme，单一真源
    "logic":   (T.BLUE_L,   T.BLUE,   T.BLUE_D),
    "memory":  (T.TEAL_L,   T.TEAL,   T.TEAL_D),
    "power":   (T.AMBER_L,  T.AMBER,  T.AMBER_D),
    "io":      (T.ROSE_L,   T.ROSE,   T.ROSE_D),
    "clock":   (T.VIOLET_L, T.VIOLET, T.VIOLET_D),
    "neutral": (T.SLATE_L,  "#94A3B8", T.INK),
}


def _drole3(role, variant):
    fill, line, tc = _DROLE.get(role, _DROLE["neutral"])
    if variant == "solid":
        return line, line, "#FFFFFF"
    if variant == "outline":
        return "#FFFFFF", line, tc
    return fill, line, tc


def _dbox(slide, x, y, w, h, text, role="neutral", variant="soft", fs=12):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.12
    except Exception:
        pass
    fc, ec, tc = _drole3(role, variant)
    shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fc)
    shp.line.color.rgb = _rgb(ec); shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3); tf.margin_top = tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        _set_font(r, size=fs, bold=True, color=tc.lstrip("#"))
    return shp


def _darrow(slide, x0, y0, x1, y1, color="#3A4252", w=1.5):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0), Inches(x1), Inches(y1))
    cn.line.color.rgb = _rgb(color.lstrip("#")); cn.line.width = Pt(w)
    cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return cn


def _dlabel(slide, x, y, text, color="#1A1F2B", fs=11, ha="center", bold=False):
    bw = 3.6
    bx = {"left": x, "center": x - bw / 2, "right": x - bw}.get(ha, x - bw / 2)
    tb = slide.shapes.add_textbox(Inches(bx), Inches(y - 0.2), Inches(bw), Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]; p.alignment = _align(ha)
    r = p.add_run(); r.text = text
    _set_font(r, size=fs, bold=bold, color=color.lstrip("#"))
    return tb


def _dia_native(slide, dia, box):
    """把一份图布局(boxes/arrows/labels，左下原点、坐标即英寸)画成 PowerPoint 原生形状。"""
    X, Ytop, BW, BH = box

    def px(x):
        return X + x

    def py(y):
        return Ytop + (BH - y)                 # 左下原点 → 幻灯片左上原点
    fs = dia.get("fs", 12)
    for x, yb, w, h, text, role, variant in dia.get("boxes", []):
        _dbox(slide, px(x), py(yb + h), w, h, text, role, variant, fs=fs)
    for a in dia.get("arrows", []):
        x0, y0, x1, y1 = a[:4]
        _darrow(slide, px(x0), py(y0), px(x1), py(y1), a[4] if len(a) > 4 else "#3A4252")
    for L in dia.get("labels", []):
        _dlabel(slide, px(L[0]), py(L[1]), L[2], L[3] if len(L) > 3 else "#1A1F2B",
                L[4] if len(L) > 4 else 11, L[5] if len(L) > 5 else "center", L[6] if len(L) > 6 else False)


def _dia_prev(ax, dia, box):
    """同一份布局的 matplotlib 预览（与原生形状一致，作为唯一的版面校对）。"""
    X, Ytop, BW, BH = box
    by = SLIDE_H - Ytop - BH

    def mx(x):
        return X + x

    def my(y):
        return by + y
    ax.add_patch(Rectangle((X, by), BW, BH, fc="none", ec="#E2E8F0", lw=1, zorder=1))
    fs = dia.get("fs", 12)
    lh = fs / 72.0 * 1.18
    for x, yb, w, h, text, role, variant in dia.get("boxes", []):
        fc, ec, tc = _drole3(role, variant)
        ax.add_patch(FancyBboxPatch((mx(x), my(yb)), w, h, boxstyle="round,pad=0,rounding_size=0.08",
                     fc=fc, ec=ec, lw=1.5, zorder=3))
        lines = []
        for seg in str(text).split("\n"):
            lines += _wrap(seg, max(1.0, (w - 0.18) * 72.0 / fs)) if seg else []
        for i, ln in enumerate(lines):
            ax.text(mx(x) + w / 2, my(yb) + h / 2 + ((len(lines) - 1) / 2 - i) * lh, ln,
                    ha="center", va="center", color=tc, fontsize=fs, fontweight="bold", zorder=4)
    for a in dia.get("arrows", []):
        x0, y0, x1, y1 = a[:4]
        ax.add_patch(FancyArrowPatch((mx(x0), my(y0)), (mx(x1), my(y1)), arrowstyle="-|>",
                     mutation_scale=11, color=a[4] if len(a) > 4 else "#3A4252", lw=1.5, zorder=4, shrinkA=2, shrinkB=2))
    for L in dia.get("labels", []):
        ax.text(mx(L[0]), my(L[1]), L[2], ha=(L[5] if len(L) > 5 else "center"), va="center",
                color=(L[3] if len(L) > 3 else "#1A1F2B"), fontsize=(L[4] if len(L) > 4 else 11),
                fontweight=("bold" if (len(L) > 6 and L[6]) else "normal"), zorder=4)


def _dia_uniquify():
    R, RD, TD, I2 = T.ROSE, T.ROSE_D, T.TEAL_D, T.INK2
    boxes = [
        (1.0, 4.0, 1.4, 0.55, "top", "neutral", "soft"),
        (0.15, 2.9, 1.3, 0.55, "m1", "logic", "soft"), (1.95, 2.9, 1.3, 0.55, "m2", "logic", "soft"),
        (1.0, 1.8, 1.4, 0.55, "amod", "memory", "soft"), (1.0, 0.7, 1.4, 0.55, "u1", "io", "soft"),
        (4.5, 4.0, 1.4, 0.55, "top", "neutral", "soft"),
        (3.65, 2.9, 1.3, 0.55, "m1", "logic", "soft"), (5.45, 2.9, 1.3, 0.55, "m2", "logic", "soft"),
        (3.65, 1.8, 1.3, 0.55, "amod1", "memory", "soft"), (5.45, 1.8, 1.3, 0.55, "amod2", "memory", "soft"),
        (3.65, 0.7, 1.3, 0.55, "u1", "io", "soft"), (5.45, 0.7, 1.3, 0.55, "u1", "io", "soft"),
    ]
    arrows = [
        (1.5, 4.0, 0.95, 3.47, I2), (1.9, 4.0, 2.45, 3.47, I2),
        (0.95, 2.9, 1.45, 2.37, R), (2.45, 2.9, 1.95, 2.37, R), (1.7, 1.8, 1.7, 1.27, I2),
        (5.0, 4.0, 4.45, 3.47, I2), (5.4, 4.0, 5.95, 3.47, I2),
        (4.3, 2.9, 4.3, 2.37, I2), (6.1, 2.9, 6.1, 2.37, I2), (4.3, 1.8, 4.3, 1.27, I2), (6.1, 1.8, 6.1, 1.27, I2),
    ]
    labels = [
        (1.7, 4.95, "非唯一：m1/m2 共享 amod", RD, 11.5, "center", True),
        (5.2, 4.95, "唯一化：各自独立副本", TD, 11.5, "center", True),
        (1.7, 0.28, "改 m1/u1 牵连 m2/u1", RD, 9.5, "center", False),
        (5.2, 0.28, "各实例可独立优化、移动", TD, 9.5, "center", False),
    ]
    return {"fs": 13, "boxes": boxes, "arrows": arrows, "labels": labels}


def _dia_hier():
    BD, TD, MU, I2 = T.BLUE_D, T.TEAL_D, T.MUTED, T.INK2
    boxes = [
        (0.15, 1.35, 2.85, 3.1, "", "neutral", "outline"),
        (0.5, 2.5, 2.15, 1.35, "整颗芯片\n一次 P&R", "logic", "solid"),
        (3.55, 1.35, 3.15, 3.1, "", "neutral", "outline"),
        (4.55, 3.5, 1.55, 0.65, "Top 集成", "neutral", "soft"),
        (3.75, 2.35, 0.85, 0.8, "Blk1", "memory", "soft"), (4.7, 2.35, 0.85, 0.8, "Blk2", "memory", "soft"),
        (5.65, 2.35, 0.85, 0.8, "Blk3", "memory", "soft"),
        (0.15, 0.2, 6.55, 0.92, "层次化省运行时间 / 内存、利于复用；但全芯片时序收敛更难，依赖引脚 / feedthrough / 时序预算 / ILM", "neutral", "soft"),
    ]
    arrows = [(4.175, 3.15, 4.9, 3.5, I2), (5.125, 3.15, 5.2, 3.5, I2), (6.075, 3.15, 5.5, 3.5, I2)]
    labels = [
        (1.575, 4.78, "扁平 Flat", BD, 14, "center", True),
        (5.125, 4.78, "层次化 Hierarchical", TD, 14, "center", True),
        (1.575, 1.62, "运行慢 / 内存大", MU, 10, "center", False),
        (4.175, 2.12, "P&R", MU, 8.5, "center", False), (5.125, 2.12, "P&R", MU, 8.5, "center", False),
        (6.075, 2.12, "P&R", MU, 8.5, "center", False),
    ]
    return {"fs": 12, "boxes": boxes, "arrows": arrows, "labels": labels}


def _dia_pnr():
    flow = [("逻辑综合\nSynth", "neutral", "outline"), ("布图规划\nFloorplan", "clock", "solid"),
            ("布局\nPlace", "neutral", "outline"), ("时钟树\nCTS", "neutral", "outline"),
            ("布线\nRoute", "neutral", "outline")]
    boxes, arrows = [], []
    fw, gap, y = 1.15, 0.27, 4.35
    for i, (t, role, var) in enumerate(flow):
        x = 0.05 + i * (fw + gap)
        boxes.append((x, y, fw, 0.72, t, role, var))
        if i > 0:
            arrows.append((x - gap, y + 0.36, x, y + 0.36, T.INK2))
    tasks = [("① die / core 几何", "logic"), ("② 宏摆放 & 朝向", "memory"),
             ("③ 电源规划 PG", "power"), ("④ 多电压 / blockage", "io")]
    for i, (t, role) in enumerate(tasks):
        boxes.append((0.1 + (i % 2) * 3.45, 2.45 - (i // 2) * 1.15, 3.2, 0.95, t, role, "solid"))
    labels = [(0.1, 5.18, "Floorplan 在 PnR 流程中的位置", T.INK, 12.5, "left", True),
              (0.1, 3.8, "Floorplan 主要任务", T.INK, 12.5, "left", True)]
    return {"fs": 12, "boxes": boxes, "arrows": arrows, "labels": labels}


def _dia_io():
    ins = ["Netlist (.v)", "LEF", "Liberty (.lib)", "SDC", "UPF", "Partition / 预算"]
    boxes = [(0.1 + (i % 2) * 1.65, 4.3 - (i // 2) * 0.78, 1.55, 0.62, t, "neutral", "soft")
             for i, t in enumerate(ins)]
    boxes.append((4.25, 3.25, 2.45, 1.3, "Floorplan\nICC2 / Innovus", "logic", "solid"))
    outs = [("DEF：宏 + PG + blockage", "memory"), ("Floorplan DB (NDM/OA)", "memory"),
            ("可布线性 / 时序初评", "io")]
    for i, (t, role) in enumerate(outs):
        boxes.append((0.1, 1.7 - i * 0.72, 6.6, 0.6, t, role, "soft"))
    arrows = [(3.4, 3.95, 4.25, 3.9, T.LINE), (5.45, 3.25, 3.5, 2.35, T.INK2)]
    labels = [(0.1, 5.0, "输入", T.INK2, 12, "left", True), (0.1, 2.45, "输出", T.INK2, 12, "left", True)]
    return {"fs": 11.5, "boxes": boxes, "arrows": arrows, "labels": labels}


def _dia_loop():
    I2 = T.INK2
    boxes = [
        (2.3, 4.0, 2.2, 0.9, "Floorplan\n定 die/macro/PG", "logic", "solid"),
        (4.7, 2.1, 2.0, 0.9, "试布局 + GR", "memory", "solid"),
        (2.3, 0.3, 2.2, 0.9, "评估\n拥塞 / 时序 / IR", "io", "solid"),
        (0.1, 2.1, 2.0, 0.9, "回退调整", "power", "solid"),
    ]
    arrows = [(4.5, 4.05, 5.45, 3.05, I2), (5.55, 2.1, 4.5, 1.05, I2),
              (2.3, 1.05, 1.35, 2.05, I2), (1.35, 3.0, 2.3, 4.05, I2)]
    labels = [(3.4, 2.55, "迭代闭环", T.MUTED, 13, "center", True),
              (3.4, 5.05, "收敛后才进详细 Placement", T.INK, 12.5, "center", True)]
    return {"fs": 12, "boxes": boxes, "arrows": arrows, "labels": labels}


DIAGRAMS = {"uniquify": _dia_uniquify, "hier": _dia_hier, "pnr": _dia_pnr, "io": _dia_io, "loop": _dia_loop}


def build_pptx(specs, out_pptx, total, author="J.C", asset_dir="", template=None, page_label=PAGE_LABEL,
               font=None, mono_font=None):
    """从规格构建可编辑 .pptx。
    font / mono_font：覆盖正文/等宽字体（默认 Microsoft YaHei / Consolas）。换机器若缺这两款字体，
    传本机可用字体名即可，避免 PowerPoint 字体替换打乱已校对的折行。构建后自动还原全局字体，
    便于同一进程内连续构建多份不同字体的 deck。"""
    validate_specs(specs)        # 缺图软检查留给 build_previews / 渲染期占位告警，避免重复刷屏
    global YAHEI, MONO
    _saved_fonts = (YAHEI, MONO)
    YAHEI = font or YAHEI
    MONO = mono_font or MONO
    try:
        return _build_pptx_impl(specs, out_pptx, total, author, asset_dir, template, page_label)
    finally:
        YAHEI, MONO = _saved_fonts


def _build_pptx_impl(specs, out_pptx, total, author="J.C", asset_dir="", template=None, page_label=PAGE_LABEL):
    """template=<.pptx 路径> 时，基于该模板（继承其母版/主题/字体），并清掉模板自带的示例幻灯片。"""
    prs = Presentation(template) if template else Presentation()
    if template:
        sld = prs.slides._sldIdLst
        for s in list(sld):
            sld.remove(s)
    prs.slide_width = Inches(SLIDE_W); prs.slide_height = Inches(SLIDE_H)
    blank = (min(prs.slide_layouts, key=lambda L: len(L.placeholders)) if template
             else prs.slide_layouts[6])
    _clean_master(prs, blank)   # 删默认 4:3 占位符/未用版式，避免页脚重叠、母版按 16:9
    _layout_chrome(prs, blank, page_label)  # 母版件下放到版式：所有内容页继承，新增页自动带样式
    fig_no = 0
    cur_ch = None
    for i, s in enumerate(specs, 1):
        sl = prs.slides.add_slide(blank)
        k = s["kind"]
        if k in ("cover", "title"):
            _cover(sl, s); continue
        if k == "close":
            _close(sl, s); continue
        if k == "section":
            _section(sl, s); cur_ch = f"第{s['num']}章 · {s['title']}"; continue
        if k == "agenda":
            _agenda(sl, s); _page(sl); continue
        _title_block(sl, s["title"], s.get("sub"), chapter=cur_ch)
        if k == "split":
            _txt(sl, (0.7, 1.72, 5.2, 5.4), _marks(s["bullets"], s.get("style", "bullet")),
                 sizes=[14] * len(s["bullets"]), colors=[INK_C] * len(s["bullets"]), space=8, line_sp=1.16, auto_fit=True)
            cap = s.get("caption")
            box = (6.0, 1.7, 6.85, 4.7 if cap else 5.25)
            if s.get("diagram"):                       # 原生框图（可编辑形状），否则插图片
                _dia_native(sl, DIAGRAMS[s["diagram"]](), box)
            else:
                _pic(sl, os.path.join(asset_dir, s["figure"]), box)
            if cap:
                fig_no += 1
                _txt(sl, (6.0, 6.48, 6.85, 0.5), [f"图 {fig_no}. {cap}"], sizes=[12], colors=[MUTED_C], align=PP_ALIGN.CENTER)
            if s.get("credit"):              # 文献插图来源标注（恰当处）
                _txt(sl, (6.0, 6.98, 6.85, 0.3), ["来源：" + s["credit"]], sizes=[9], colors=[MUTED_C], align=PP_ALIGN.RIGHT)
        elif k == "table":
            _table(sl, s)
        elif k == "chart":
            _chart(sl, s)
        elif k == "refs":
            _refs(sl, s)
        elif k == "cards2":
            for j, (hdr, code, a) in enumerate(s["cards"]):
                x = 0.7 + j * 6.1
                _card(sl, x, 1.7, 5.8, 5.1, accent=a)
                _txt(sl, (x + 0.35, 1.95, 5.2, 0.5), [hdr], sizes=[15], bolds=[True], colors=[a])
                _txt(sl, (x + 0.35, 2.6, 5.3, 4.0), code, sizes=[11.5] * len(code), monos=[True] * len(code),
                     colors=[INK_C] * len(code), space=7, line_sp=1.12, auto_fit=True)
        elif k == "grid6":
            for j, (t1, t2, c) in enumerate(s["items"]):
                x = 0.7 + (j % 3) * 4.05
                y = 1.75 + (j // 3) * 2.45
                _card(sl, x, y, 3.75, 2.15, accent=c)
                _txt(sl, (x + 0.3, y + 0.28, 3.3, 0.6), [t1], sizes=[15], bolds=[True], colors=[INK_C])
                _txt(sl, (x + 0.3, y + 1.05, 3.3, 0.9), ["→ " + t2], sizes=[12], colors=[MUTED_C])
        elif k == "cols2":
            for j, (hdr, items, a) in enumerate(s["cols"]):
                x = 0.7 + j * 6.1
                _card(sl, x, 1.7, 5.8, 5.1, accent=a)
                _txt(sl, (x + 0.35, 1.95, 5.2, 0.5), [hdr], sizes=[15], bolds=[True], colors=[a])
                _txt(sl, (x + 0.35, 2.55, 5.3, 4.2), _bullets_lines(items),
                     sizes=[12.5] * len(items), colors=[INK_C] * len(items), space=6, line_sp=1.16, auto_fit=True)
        elif k == "bullets":
            items = s["bullets"]
            if s.get("two_col"):
                mid = (len(items) + 1) // 2
                _txt(sl, (0.7, 1.75, 6.0, 5.3), _bullets_lines(items[:mid]),
                     sizes=[14.5] * mid, colors=[INK_C] * mid, space=9, line_sp=1.2, auto_fit=True)
                _txt(sl, (6.9, 1.75, 6.0, 5.3), _bullets_lines(items[mid:]),
                     sizes=[14.5] * (len(items) - mid), colors=[INK_C] * (len(items) - mid), space=9, line_sp=1.2, auto_fit=True)
            else:
                _txt(sl, (0.7, 1.75, 12.0, 5.3), _bullets_lines(items),
                     sizes=[15] * len(items), colors=[INK_C] * len(items), space=10, line_sp=1.25, auto_fit=True)
        _page(sl)
    try:
        prs.core_properties.author = author
    except Exception as e:
        _log.debug("设置 author 失败: %s", e)
    out_dir = os.path.dirname(out_pptx)
    if out_dir:                                  # 裸文件名时 dirname 为 ''，makedirs('') 会报错
        os.makedirs(out_dir, exist_ok=True)
    prs.save(out_pptx)
    return out_pptx


# ============================ 预览端 (matplotlib) ========================= #
def _mtext(ax, x, ytop, s, fs=14, color=T.INK, bold=False, mono=False, ha="left", va="top"):
    ax.text(x, SLIDE_H - ytop, s, fontsize=fs, color=color, ha=ha, va=va,
            fontweight=("bold" if bold else "normal"),
            family=("monospace" if mono else "sans-serif"))


def _mcard(ax, x, ytop, w, h, fc="#FFFFFF", ec=T.LINE, accent=None):
    ax.add_patch(FancyBboxPatch((x, SLIDE_H - ytop - h), w, h, boxstyle="round,pad=0,rounding_size=0.08",
                 fc=fc, ec=ec, lw=1.3, zorder=2))
    if accent:
        ax.add_patch(Rectangle((x, SLIDE_H - ytop - h + 0.12), 0.09, h - 0.24, fc=accent, ec="none", zorder=3))


def _mmissing(ax, path, box):
    """预览端缺图占位：灰色虚线框 + 居中红字，与 pptx 端 _missing_pic 一致。"""
    x, t, w, h = box
    ax.add_patch(FancyBboxPatch((x, SLIDE_H - t - h), w, h, boxstyle="round,pad=0,rounding_size=0.08",
                 fc="#EDF0F4", ec="#C7CDD6", lw=1.3, ls="--", zorder=4))
    ax.text(x + w / 2, SLIDE_H - t - h / 2, "缺图 missing:\n" + os.path.basename(path),
            ha="center", va="center", color="#A6473C", fontsize=13, fontweight="bold", zorder=5)


def _wrap(s, maxu):
    """按显示宽度折行（CJK≈1，ASCII≈0.55 单位），让预览贴近 pptx 的自动换行。
    连续 ASCII 字母/数字成"词"整体折行，不拆词（与 PowerPoint 的按词换行一致）。"""
    toks, i, n = [], 0, len(s)
    while i < n:
        if s[i].isascii() and s[i].isalnum():
            j = i + 1
            while j < n and s[j].isascii() and s[j].isalnum():
                j += 1
            toks.append(s[i:j]); i = j
        else:
            toks.append(s[i]); i += 1
    out, line, u = [], "", 0.0
    for t in toks:
        w = sum(T.ASCII_RATIO if ord(c) < 128 else 1.0 for c in t)
        if u + w > maxu and line:
            out.append(line.rstrip()); line, u = t.lstrip(), w
        else:
            line += t; u += w
    if line.strip():
        out.append(line.rstrip())
    return out or [""]


def _maxu(width_in, fs):
    """正文框可用宽度(英寸) → _wrap 的 maxu（≈可容纳的 CJK 字数）。
    供新页型按真实框宽推折行阈值，取代经验魔数；现有页型沿用其手调常量。"""
    return width_in / (fs / 72.0)


def _pbul(ax, xb, xt, yy, text, fs, color, maxu, lh=0.40, gap=0.16, mark="▪"):
    """预览端：画一条带标记(▪ 或 ①)的可折行要点，返回下一行 y。"""
    _mtext(ax, xb, yy, mark, fs=fs, color=color, bold=True)
    for ln in _wrap(text, maxu):
        _mtext(ax, xt, yy, ln, fs=fs)
        yy += lh
    return yy + gap


def _fit_scale(items, maxu, top, bottom, lh, gap, floor=0.60):
    """模拟 PowerPoint 的 shrink-to-fit：估算这组要点在 base 行高下的总高度，
    超出可用高度则返回 <1 的缩放，使预览端按相同比例缩字号 → 预览≈导出 pptx。
    （pptx 端 auto_fit=True 让 PowerPoint 自己缩；此处让预览跟着缩，校对才准。）"""
    total = sum(len(_wrap(it, maxu)) * lh + gap for it in items)
    avail = bottom - top
    if total <= avail or total <= 0:
        return 1.0
    return max(floor, avail / total)


def _pcover(ax, s):
    ax.add_patch(Rectangle((0, 0), SLIDE_W, SLIDE_H, fc="white", ec="none", zorder=0))
    ax.add_patch(Rectangle((0, 0), 0.30, SLIDE_H, fc="#" + PRIMARY, ec="none", zorder=2))
    _mtext(ax, 1.15, 1.45, s.get("tag", "数字 IC 后端 · DVD Lecture 6"), fs=13, bold=True, color="#" + PRIMARY)
    _mtext(ax, 1.1, 2.45, s["title"], fs=40, bold=True, color="#" + INK_C)
    if s.get("sub"):
        _mtext(ax, 1.15, 3.95, s["sub"], fs=17, color="#" + SUB_C)
    ax.add_patch(Rectangle((1.15, SLIDE_H - 4.80 - 0.05), 3.8, 0.05, fc="#" + PRIMARY, ec="none", zorder=2))
    if s.get("line"):
        _mtext(ax, 1.15, 5.18, s["line"], fs=13, bold=True, color="#" + ACCENT)
    _mtext(ax, 1.15, 6.9, s.get("src", ""), fs=10, color="#" + MUTED_C)


def _pclose(ax, s):
    ax.add_patch(Rectangle((0, 0), SLIDE_W, SLIDE_H, fc="white", ec="none", zorder=0))
    ax.add_patch(Rectangle((0, 0), 0.30, SLIDE_H, fc="#" + PRIMARY, ec="none", zorder=2))
    _mtext(ax, 1.1, 2.75, s["title"], fs=38, bold=True, color="#" + INK_C)
    if s.get("sub"):
        _mtext(ax, 1.15, 4.15, s["sub"], fs=16, color="#" + SUB_C)
    ax.add_patch(Rectangle((1.15, SLIDE_H - 4.95 - 0.05), 3.4, 0.05, fc="#" + ACCENT, ec="none", zorder=2))
    if s.get("line"):
        _mtext(ax, 1.15, 5.40, s["line"], fs=12, color="#" + MUTED_C)
    _mtext(ax, 1.15, 6.9, s.get("src", ""), fs=10, color="#" + MUTED_C)


def _psection(ax, s):
    ax.add_patch(Rectangle((0, 0), SLIDE_W, SLIDE_H, fc="white", ec="none", zorder=0))
    pw = 4.5
    ax.add_patch(Rectangle((0, 0), pw, SLIDE_H, fc="#" + PRIMARY, ec="none", zorder=1))
    _mtext(ax, pw / 2, 1.65, s.get("eyebrow", "CHAPTER"), fs=15, bold=True, color="#C5D6F7", ha="center")
    _mtext(ax, pw / 2, 3.6, str(s["num"]), fs=135, bold=True, color="white", ha="center", va="center")
    _mtext(ax, 5.1, 1.95, s["title"], fs=32, bold=True, color="#" + INK_C)
    if s.get("sub"):
        _mtext(ax, 5.15, 3.2, s["sub"], fs=15, color="#" + SUB_C)
    ax.add_patch(Rectangle((5.15, SLIDE_H - 3.95 - 0.05), 2.4, 0.05, fc="#" + ACCENT, ec="none", zorder=2))
    items = s.get("items", []); dy = min(0.52, 2.9 / max(1, len(items)))
    for i, it in enumerate(items):
        y = 4.35 + i * dy
        ax.add_patch(Rectangle((5.15, SLIDE_H - y - 0.19), 0.13, 0.13, fc="#" + ACCENT, ec="none", zorder=2))
        _mtext(ax, 5.45, y, it, fs=13, color="#" + INK_C)


def _pagenda(ax, s, i, page_label):
    ax.add_patch(Rectangle((0, 0), SLIDE_W, SLIDE_H, fc="white", ec="none", zorder=0))
    ax.add_patch(Rectangle((0.6, SLIDE_H - 0.5 - 0.62), 0.16, 0.62, fc="#" + PRIMARY, ec="none", zorder=2))
    _mtext(ax, 0.85, 0.5, s["title"], fs=24, bold=True)
    if s.get("sub"):
        _mtext(ax, 0.87, 1.2, s["sub"], fs=13, color=T.MUTED)
    ax.add_patch(Rectangle((0.85, SLIDE_H - 1.6 - 0.02), 11.85, 0.02, fc="#" + HAIR_C, ec="none", zorder=1))
    secs = s["sections"]; colx = [0.95, 7.0]
    n = len(secs); per = max(1, (n + 1) // 2); top0 = 2.05; rowh = min(1.16, 4.15 / per)
    for j, (num, ttl, det) in enumerate(secs):
        x = colx[j // per]; y = top0 + (j % per) * rowh
        ax.add_patch(Circle((x + 0.23, SLIDE_H - y - 0.23), 0.23, fc="#" + PRIMARY, ec="none", zorder=3))
        _mtext(ax, x + 0.23, y + 0.23, str(num), fs=14, bold=True, color="white", ha="center", va="center")
        _mtext(ax, x + 0.64, y - 0.02, ttl, fs=15, bold=True, color="#" + INK_C)
        _mtext(ax, x + 0.64, y + 0.42, det, fs=12, color="#" + MUTED_C)
    if s.get("line"):
        ax.add_patch(FancyBboxPatch((0.95, SLIDE_H - 6.28 - 0.6), 11.0, 0.6,
                     boxstyle="round,pad=0,rounding_size=0.06", fc="#F6E8D5", ec="#E4C79A", lw=1.2, zorder=2))
        _mtext(ax, 1.25, 6.5, s["line"], fs=12.5, bold=True, color="#8A5212")
    _mtext(ax, 0.6, 7.1, page_label, fs=10, color="#" + PAGE_C)
    _mtext(ax, 12.9, 7.1, str(i), fs=10, color="#" + PAGE_C, ha="right")


def _ptable(ax, s):
    t = s["table"]; headers = t["headers"]; rows = t["rows"]
    aligns = t.get("col_align"); widths = t.get("col_w")
    ncol = len(headers); nrow = len(rows) + 1
    x, top, w, h = 0.7, 1.78, 11.95, 5.05
    cw = [w * v / float(sum(widths)) for v in widths] if widths else [w / ncol] * ncol
    xstops = [x]
    for c in range(ncol):
        xstops.append(xstops[-1] + cw[c])
    rh = h / nrow
    al = lambda c: (aligns[c] if aligns and c < len(aligns) else "left")

    def cell(rr, cc, text, fill, color, bold):
        cx0 = xstops[cc]; cyt = top + rr * rh
        ax.add_patch(Rectangle((cx0, SLIDE_H - cyt - rh), cw[cc], rh, fc=fill, ec="#C7CDD6", lw=0.8, zorder=2))
        ha = al(cc); tx = cx0 + 0.12 if ha == "left" else (cx0 + cw[cc] - 0.12 if ha == "right" else cx0 + cw[cc] / 2)
        ha2 = {"left": "left", "right": "right", "center": "center"}.get(ha, "left")
        ax.text(tx, SLIDE_H - cyt - rh / 2, text, ha=ha2, va="center", color=color,
                fontsize=10.5 if bold else 10, fontweight=("bold" if bold else "normal"), zorder=3)

    for c, htext in enumerate(headers):
        cell(0, c, str(htext), "#" + PRIMARY, "white", True)
    for ri, row in enumerate(rows, 1):
        fill = "#FFFFFF" if ri % 2 else "#EDF0F4"
        for c in range(ncol):
            cell(ri, c, str(row[c]) if c < len(row) else "", fill, "#" + INK_C, False)


def _pchart(ax, s):
    c = s["chart"]; ctype = c.get("type", "line")
    x, top, w, h = 0.7, 1.85, 11.95, 4.85
    iax = ax.inset_axes([x, SLIDE_H - top - h, w, h], transform=ax.transData)
    cats = c["categories"]; xs = list(range(len(cats)))
    pal = ["#" + c for c in SERIES_COLORS]
    if ctype == "line":
        for i, (name, vals) in enumerate(c["series"]):
            iax.plot(xs, vals, marker="o", ms=4, lw=2, color=pal[i % len(pal)], label=name)
    else:
        n = len(c["series"]); bw = 0.8 / max(n, 1)
        for i, (name, vals) in enumerate(c["series"]):
            iax.bar([xx + (i - (n - 1) / 2.0) * bw for xx in xs], vals, width=bw, color=pal[i % len(pal)], label=name)
    iax.set_xticks(xs); iax.set_xticklabels(cats, fontsize=9)
    iax.tick_params(labelsize=8)
    if len(c["series"]) > 1:
        iax.legend(fontsize=8, loc="best")
    if c.get("x_title"):
        iax.set_xlabel(c["x_title"], fontsize=9)
    if c.get("y_title"):
        iax.set_ylabel(c["y_title"], fontsize=9)


def _oflow_mark(ax):
    """预览端唯一能本地算出的版面警告：正文超出可用高度时，在底部中央打红字告警。
    pptx 端 PowerPoint 会 shrink-to-fit 自救，但预览是唯一校对通道，必须把溢出显式标出来。"""
    _mtext(ax, SLIDE_W / 2.0, 7.30, "⚠ OVERFLOW 正文超出版面（请精简内容或缩小字号）",
           fs=11, color="#C0152F", bold=True, ha="center")


def _prefs(ax, s):
    items = s["refs"]
    mx = 1.95
    if len(items) > 8:
        mid = (len(items) + 1) // 2
        yy = 1.95
        for i, r in enumerate(items[:mid]):
            yy = _pbul(ax, 0.7, 1.15, yy, r, 12, "#" + INK_C, 26, mark=f"[{i + 1}]")
        mx = yy; yy = 1.95
        for i, r in enumerate(items[mid:]):
            yy = _pbul(ax, 6.9, 7.35, yy, r, 12, "#" + INK_C, 26, mark=f"[{mid + i + 1}]")
        mx = max(mx, yy)
    else:
        yy = 1.95
        for i, r in enumerate(items):
            yy = _pbul(ax, 0.7, 1.15, yy, r, 12.5, "#" + INK_C, 52, mark=f"[{i + 1}]")
        mx = yy
    return mx


def build_previews(specs, outdir, total, asset_dir="", page_label=PAGE_LABEL):
    os.makedirs(outdir, exist_ok=True)
    validate_specs(specs, asset_dir)
    paths = []
    T.setup_fonts()
    fig_no = 0
    cur_ch = None
    for i, s in enumerate(specs, 1):
        fig, ax = plt.subplots(figsize=(SLIDE_W, SLIDE_H), dpi=110)
        ax.set_xlim(0, SLIDE_W); ax.set_ylim(0, SLIDE_H); ax.axis("off"); ax.set_autoscale_on(False)
        fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
        k = s["kind"]
        mark_c = "#" + INK_C    # 要点标记色：与 pptx 一致（pptx 整行 INK），不用 accent，避免预览与 pptx 不符
        if k in ("cover", "title"):
            _pcover(ax, s); paths.append(_save_prev(fig, outdir, i)); continue
        if k == "close":
            _pclose(ax, s); paths.append(_save_prev(fig, outdir, i)); continue
        if k == "section":
            _psection(ax, s); cur_ch = f"第{s['num']}章 · {s['title']}"; paths.append(_save_prev(fig, outdir, i)); continue
        if k == "agenda":
            _pagenda(ax, s, i, page_label); paths.append(_save_prev(fig, outdir, i)); continue
        ax.add_patch(Rectangle((0, 0), SLIDE_W, SLIDE_H, fc="white", ec="none", zorder=0))
        ax.add_patch(Rectangle((0.6, SLIDE_H - 0.5 - 0.62), 0.16, 0.62, fc="#" + PRIMARY, ec="none", zorder=2))
        _mtext(ax, 0.85, 0.5, s["title"], fs=24, bold=True)
        if s.get("sub"):
            _mtext(ax, 0.87, 1.2, s["sub"], fs=13, color=T.MUTED)
        if cur_ch:
            _mtext(ax, 12.7, 0.62, cur_ch, fs=11, bold=True, color="#" + PRIMARY, ha="right")
        ax.add_patch(Rectangle((0.85, SLIDE_H - 1.6 - 0.02), 11.85, 0.02, fc="#" + HAIR_C, ec="none", zorder=1))
        _mtext(ax, 0.6, 7.1, page_label, fs=10, color="#" + PAGE_C)
        if k == "split":
            style = s.get("style", "bullet")
            maxu = 26 if style == "prose" else 24
            blh, bgap = (0.40, 0.22) if style == "prose" else (0.40, 0.16)
            sc = _fit_scale(s["bullets"], maxu, 1.95, 6.98, blh, bgap)   # 模拟 pptx 自动缩字号
            yy = 1.95
            for bi, b in enumerate(s["bullets"]):
                if style == "prose":                       # 叙述性内容：整段、无标记、段间留白
                    yy = _pbul(ax, 0.7, 0.7, yy, b, 14 * sc, mark_c, maxu, lh=blh * sc, gap=bgap * sc, mark="")
                else:
                    mk = _circ(bi) if style == "num" else "▪"
                    yy = _pbul(ax, 0.7, 1.05, yy, b, 14 * sc, mark_c, maxu, lh=blh * sc, gap=bgap * sc, mark=mk)
            if yy > 7.05:                                   # 缩到底仍溢出才告警
                _oflow_mark(ax)
            cap = s.get("caption")
            box = (6.0, 1.7, 6.85, 4.7 if cap else 5.25)
            if s.get("diagram"):
                _dia_prev(ax, DIAGRAMS[s["diagram"]](), box)
            else:
                p = os.path.join(asset_dir, s["figure"])
                if not os.path.isfile(p):
                    _mmissing(ax, p, box)
                else:
                    iw, ih = Image.open(p).size
                    fx, ft, fw, fh = fit(box, iw, ih)
                    ax.imshow(mpimg.imread(p), extent=[fx, fx + fw, SLIDE_H - ft - fh, SLIDE_H - ft], aspect="auto", zorder=4)
                    ax.add_patch(Rectangle((box[0], SLIDE_H - box[1] - box[3]), box[2], box[3], fc="none", ec="#E2E8F0", lw=1, zorder=1))
            if cap:
                fig_no += 1
                _mtext(ax, 6.0 + 6.85 / 2.0, 6.62, f"图 {fig_no}. {cap}", fs=11, color="#" + MUTED_C, ha="center")
            if s.get("credit"):
                _mtext(ax, 12.85, 7.02, "来源：" + s["credit"], fs=9, color="#" + MUTED_C, ha="right")
        elif k == "table":
            _ptable(ax, s)
        elif k == "chart":
            _pchart(ax, s)
        elif k == "refs":
            if _prefs(ax, s) > 6.98:
                _oflow_mark(ax)
        elif k == "cards2":
            mx = 0
            for j, (hdr, code, a) in enumerate(s["cards"]):
                x = 0.7 + j * 6.1
                _mcard(ax, x, 1.7, 5.8, 5.1, accent="#" + a if not a.startswith("#") else a)
                _mtext(ax, x + 0.35, 1.95, hdr, fs=14, bold=True, color="#" + a if not a.startswith("#") else a)
                yy = 2.6
                for ln in code:
                    _mtext(ax, x + 0.35, yy, ln, fs=11, mono=True); yy += 0.42
                mx = max(mx, yy)
            if mx > 6.85:
                _oflow_mark(ax)
        elif k == "grid6":
            for j, (t1, t2, c) in enumerate(s["items"]):
                x = 0.7 + (j % 3) * 4.05; y = 1.75 + (j // 3) * 2.45
                _mcard(ax, x, y, 3.75, 2.15, accent=c if c.startswith("#") else "#" + c)
                _mtext(ax, x + 0.3, y + 0.28, t1, fs=14, bold=True)
                _mtext(ax, x + 0.3, y + 1.05, "→ " + t2, fs=12, color=T.MUTED)
        elif k == "cols2":
            mx = 0
            for j, (hdr, items, a) in enumerate(s["cols"]):
                x = 0.7 + j * 6.1
                ac = a if a.startswith("#") else "#" + a
                _mcard(ax, x, 1.7, 5.8, 5.1, accent=ac)
                _mtext(ax, x + 0.35, 1.95, hdr, fs=14, bold=True, color=ac)
                sc = _fit_scale(items, 28, 2.55, 6.72, 0.36, 0.12)
                yy = 2.55
                for it in items:
                    yy = _pbul(ax, x + 0.32, x + 0.6, yy, it, 12.5 * sc, mark_c, 28, lh=0.36 * sc, gap=0.12 * sc)
                mx = max(mx, yy)
            if mx > 6.8:
                _oflow_mark(ax)
        elif k == "bullets":
            items = s["bullets"]
            if s.get("two_col"):
                mid = (len(items) + 1) // 2
                scL = _fit_scale(items[:mid], 30, 1.95, 6.98, 0.40, 0.16)
                scR = _fit_scale(items[mid:], 30, 1.95, 6.98, 0.40, 0.16)
                yy = 1.95
                for b in items[:mid]:
                    yy = _pbul(ax, 0.7, 1.0, yy, b, 14 * scL, mark_c, 30, lh=0.40 * scL, gap=0.16 * scL)
                ya = yy; yy = 1.95
                for b in items[mid:]:
                    yy = _pbul(ax, 6.9, 7.2, yy, b, 14 * scR, mark_c, 30, lh=0.40 * scR, gap=0.16 * scR)
                if max(ya, yy) > 7.05:
                    _oflow_mark(ax)
            else:
                sc = _fit_scale(items, 52, 1.95, 6.98, 0.40, 0.16)
                yy = 1.95
                for b in items:
                    yy = _pbul(ax, 0.7, 1.0, yy, b, 15 * sc, mark_c, 52, lh=0.40 * sc, gap=0.16 * sc)
                if yy > 7.05:
                    _oflow_mark(ax)
        _mtext(ax, 12.9, 7.1, str(i), fs=10, color="#" + PAGE_C, ha="right")
        paths.append(_save_prev(fig, outdir, i))
    return paths


def _save_prev(fig, outdir, i):
    p = os.path.join(outdir, f"prev_{i:02d}.png")
    fig.savefig(p, dpi=110); plt.close(fig)
    return p
