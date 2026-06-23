# -*- coding: utf-8 -*-
"""Dump a .pptx 的母版主题：尺寸、颜色方案、字体、母版/版式名。用法: python tools/inspect_pptx.py <file>"""
import sys
import zipfile
import re
from pptx import Presentation

f = sys.argv[1]
p = Presentation(f)
print("slide size:", round(p.slide_width / 914400, 3), "x", round(p.slide_height / 914400, 3), "in")
print("masters:", len(p.slide_masters), "| layouts:", len(p.slide_layouts))
for i, lay in enumerate(p.slide_layouts):
    ph = [s.name for s in lay.placeholders]
    print(f"  layout[{i}] '{lay.name}' placeholders={ph}")

z = zipfile.ZipFile(f)
theme = [n for n in z.namelist() if re.match(r"ppt/theme/theme\d+\.xml", n)]
xml = z.read(theme[0]).decode("utf-8", "ignore")
clr = z.read(theme[0]).decode("utf-8", "ignore")
scheme = re.search(r"<a:clrScheme.*?</a:clrScheme>", xml, re.S)
if scheme:
    pairs = re.findall(r'<a:(\w+)>\s*<a:(?:srgbClr val="([0-9A-Fa-f]{6})"|sysClr[^>]*lastClr="([0-9A-Fa-f]{6})")', scheme.group(0))
    print("color scheme:")
    for name, a, b in pairs:
        print(f"    {name:10s} #{(a or b).upper()}")
fonts = re.findall(r'<a:(major|minor)Font>\s*<a:latin typeface="([^"]*)"(?:[^>]*)?/>', xml)
print("fonts:", fonts)
ea = re.findall(r'<a:(major|minor)Font>.*?<a:ea typeface="([^"]*)"', xml, re.S)
print("ea fonts:", ea)
